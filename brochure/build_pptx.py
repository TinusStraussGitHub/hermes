#!/usr/bin/env python3
"""Build a BME-branded BLASTMAP UNDERGROUND brochure PowerPoint - images matched by original filename."""
import sys
from pathlib import Path

sys.path.insert(0, '/opt/data/venv/lib/python3.13/site-packages')
sys.path.insert(0, '/opt/data/skills/productivity/bme-powerpoint/scripts')
import bme_pptx_kit as K
from pptx.util import Inches, Pt

HERE = Path('/opt/data/home/hermes/brochure')
prs = K.create_presentation()

# Original filenames tell us what each image shows:
# img_2 Vibration Prediction, img_3 Fragmentation Prediction,
# img_4 Shaft Design3D, img_5 Shaft Design,
# img_6 Ring Design3D, img_7 Working Plan Report, img_8 Technical Report,
# img_9 Ring Design, img_10 Tunnel Design3D, img_11 Tunnel Design,
# img_12 Cut Designer, img_13 Shape Designer
IMG = {n: str(HERE/f'brochure_img_{n}.png') for n in range(2,14)}

# Title slide
K.add_title_slide(prs, "BLASTMAP UNDERGROUND",
                  "Complete Underground Drill & Blast Design Solution")
slides = list(prs.slides)
tb = slides[0].shapes.add_textbox(Inches(0.6), Inches(4.4), Inches(12), Inches(1.0))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Design Smarter. Blast Better. Deliver Results."
p.font.size = Pt(20); p.font.color.rgb = K.BME_RED; p.font.bold = True

# Intro slide
s = K.blank_slide(prs)
K.add_logo(s)
K.add_title_bar(s, "Underground Blast Design, Reimagined")
box = K.content_box(s, left=Inches(0.6), top=Inches(2.0), width=Inches(12.1), height=Inches(3.2))
tf = box.text_frame; tf.word_wrap = True
tf.text = ("BLASTMAP UNDERGROUND is a powerful mine planning and blast design solution built "
           "specifically for underground mining operations. From development ends and production "
           "stopes to shafts and longhole drilling, BLASTMAP UNDERGROUND provides the tools "
           "engineers need to design, analyse, and optimize every blast.")
tf.paragraphs[0].font.size = Pt(18); tf.paragraphs[0].font.color.rgb = K.BME_BLACK
K.add_footer(s, "BLASTMAP UNDERGROUND  |  Member of the Omnia Group", "www.bme.co.za")

def add_feature_slide(name, desc, img):
    s = K.blank_slide(prs)
    K.add_logo(s)
    K.add_title_bar(s, name)
    box = K.content_box(s, left=Inches(0.6), top=Inches(2.0), width=Inches(7.4), height=Inches(4.4))
    tf = box.text_frame; tf.word_wrap = True
    tf.text = desc
    tf.paragraphs[0].font.size = Pt(16); tf.paragraphs[0].font.color.rgb = K.BME_BLACK
    K.add_image(s, img, left=Inches(8.3), top=Inches(2.1), max_w=Inches(4.4), max_h=Inches(4.2))
    K.add_footer(s, "BLASTMAP UNDERGROUND", "www.bme.co.za")

# Key Features - matched by filename
add_feature_slide("Tunnel Shape Designer",
    "Quickly create and customize tunnel profiles to match your mine's specific development "
    "requirements. Design accurate tunnel geometries with confidence and ensure consistent "
    "blast layouts throughout your operation.", IMG[11])  # Tunnel Design.png

add_feature_slide("Reusable Cut Designer",
    "Save time and standardize designs with the Cut Designer. Create, store, and reuse proven "
    "cut patterns, allowing engineers to apply best-practice designs across multiple headings "
    "and projects.", IMG[12])  # Cut Designer.png

K.add_section_slide(prs, "Shaft & Ring Design", "Built for underground geometry")
add_feature_slide("Shaft Design with Rings",
    "Design shaft blasting patterns efficiently using ring-based layouts. Generate accurate "
    "shaft blast designs that improve planning consistency and support safe, productive "
    "operations.", IMG[5])  # Shaft Design.png
add_feature_slide("Uphole Ring Design",
    "Develop and optimize uphole ring patterns for production blasting. Easily manage drilling "
    "parameters and ring layouts to maximize blast performance and ore recovery.", IMG[9])  # Ring Design.png

# Reporting
K.add_section_slide(prs, "Powerful Reporting", "Professional output for every blast")
add_feature_slide("Technical Reports",
    "Generate comprehensive technical reports containing detailed blast design information, "
    "drilling parameters, and engineering data for review and auditing purposes.", IMG[8])  # Technical Report.png
add_feature_slide("Working Plan Reports",
    "Produce clear, practical working plans for field crews, ensuring accurate communication "
    "of blast instructions and improving execution underground.", IMG[7])  # Working Plan Report.png

# Advanced Analysis
K.add_section_slide(prs, "Advanced Blast Analysis", "Predict before you fire")
add_feature_slide("Vibration Prediction Model",
    "Predict blast-induced ground vibrations before firing. Evaluate potential impacts and make "
    "informed design adjustments to help maintain compliance and reduce operational risk.", IMG[2])  # Vibration Prediction.png
add_feature_slide("Fragmentation Prediction Model",
    "Estimate expected rock fragmentation outcomes and optimize blast designs to improve "
    "downstream performance, material handling, and overall productivity.", IMG[3])  # Fragmentation Prediction.png

# Why choose
s = K.blank_slide(prs)
K.add_logo(s)
K.add_title_bar(s, "Why Choose BLASTMAP UNDERGROUND?")
K.add_bullets(s, [
    "Faster blast design workflows",
    "Standardized and reusable design templates",
    "Improved blast quality and consistency",
    "Advanced predictive analysis tools",
    "Comprehensive technical and operational reporting",
    "Designed specifically for underground mining environments",
], left=Inches(0.8), top=Inches(2.0), width=Inches(11.8))
K.add_footer(s, "BLASTMAP UNDERGROUND", "www.bme.co.za")

# Closing
s = K.blank_slide(prs)
K.add_logo(s)
K.add_title_bar(s, "Transform Your Underground Blast Planning")
box = K.content_box(s, left=Inches(0.6), top=Inches(2.0), width=Inches(12.1), height=Inches(3.2))
tf = box.text_frame; tf.word_wrap = True
tf.text = ("BLASTMAP UNDERGROUND combines powerful design tools, advanced analysis capabilities, "
           "and professional reporting into a single integrated platform—helping mining operations "
           "achieve safer, more efficient, and more predictable blasting results.")
tf.paragraphs[0].font.size = Pt(18); tf.paragraphs[0].font.color.rgb = K.BME_BLACK
K.add_footer(s, "BLASTMAP UNDERGROUND  |  Member of the Omnia Group", "www.bme.co.za")

out = HERE / 'BLASTMAP_UNDERGROUND_Brochure.pptx'
K.save(prs, str(out))
print("PPTX written:", out, "slides:", len(prs.slides._sldIdLst))
